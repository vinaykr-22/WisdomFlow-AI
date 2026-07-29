from pathlib import Path
from typing import Any
from fastapi import HTTPException
import pypdf
from docx import Document as DocxDocument


def resolve_file_path(file_path: str) -> Path:
    if not file_path:
        raise HTTPException(
            status_code=404,
            detail="Document file path is missing or invalid. Please re-upload the document."
        )

    path = Path(file_path)
    if path.is_file():
        return path

    filename = path.name

    # 1. Try relative to backend root directory / uploads
    backend_dir = Path(__file__).resolve().parent.parent.parent
    local_path = backend_dir / "uploads" / filename
    if local_path.is_file():
        return local_path

    # 2. Try current working directory / uploads
    cwd_path = Path.cwd() / "uploads" / filename
    if cwd_path.is_file():
        return cwd_path

    # 3. Try current working directory directly
    cwd_filename = Path.cwd() / filename
    if cwd_filename.is_file():
        return cwd_filename

    raise HTTPException(
        status_code=404,
        detail=f"Document file '{filename}' was not found on server storage. Please re-upload the document."
    )


def extract_text(file_path: str) -> str:
    path = resolve_file_path(file_path)
    ext = path.suffix.lower()

    if ext == ".pdf":
        text = ""
        with open(path, "rb") as f:
            reader = pypdf.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() or ""
        return text

    if ext == ".docx":
        doc = DocxDocument(path)
        return "\n".join(p.text for p in doc.paragraphs)

    if ext == ".txt":
        return path.read_text(encoding="utf-8", errors="replace")

    if ext in (".ppt", ".pptx"):
        from pptx import Presentation
        prs = Presentation(path)
        return "\n".join(
            shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
        )

    raise ValueError(f"Unsupported file type: {ext}")


def extract_images(file_path: str) -> list[dict[str, Any]]:
    path = resolve_file_path(file_path)
    ext = path.suffix.lower()
    images: list[dict[str, Any]] = []

    if ext == ".pdf":
        with open(path, "rb") as f:
            reader = pypdf.PdfReader(f)
            for i, page in enumerate(reader.pages):
                for img in page.images:
                    images.append({
                        "name": img.name or f"page_{i + 1}_{len(images)}.png",
                        "data": img.data,
                        "content_type": "image/png",
                        "page": i + 1,
                    })

    elif ext == ".docx":
        doc = DocxDocument(path)
        for rel in doc.part.rels.values():
            if "image" in rel.reltype and hasattr(rel, "target_part") and rel.target_part:
                images.append({
                    "name": rel.target_ref.split("/")[-1],
                    "data": rel.target_part.blob,
                    "content_type": rel.target_part.content_type or "image/png",
                    "page": 1,
                })

    elif ext in (".ppt", ".pptx"):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(path)
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    images.append({
                        "name": shape.image.content_type.split("/")[-1] or f"slide_{slide_num}_{len(images)}.png",
                        "data": shape.image.blob,
                        "content_type": shape.image.content_type,
                        "page": slide_num,
                    })

    return images

